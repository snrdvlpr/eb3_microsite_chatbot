"""
PowerPoint (PPTX) text extraction for slide decks.
"""
from __future__ import annotations


def extract_text_from_pptx_bytes(data: bytes) -> str:
    """
    Extract visible text from a PPTX file given as bytes.
    Only pulls text from shapes; image-only slides will appear empty and may require OCR.
    """
    import io
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                t = text.strip()
                if t:
                    parts.append(t)

    return "\n\n".join(parts).strip()

