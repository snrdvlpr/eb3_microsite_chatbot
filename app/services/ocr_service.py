"""
OCR service: extract text from image-based documents (PDF, PPTX, images) using Tesseract.

This module is intentionally self-contained so you can swap out the OCR backend
later (e.g. Google Vision, AWS Textract) without touching ingestion_service.
"""
from __future__ import annotations

from typing import Optional

from app.utils.text_cleaning import clean_extracted_text


def _ocr_image(image) -> str:
    """Run Tesseract OCR on a single PIL.Image with basic preprocessing."""
    import pytesseract
    from PIL import Image

    # Ensure RGB then upscale small images to help OCR quality
    if not isinstance(image, Image.Image):
        image = Image.fromarray(image)
    if image.mode not in ("L", "RGB"):
        image = image.convert("RGB")
    min_dim = 800
    if image.width < min_dim or image.height < min_dim:
        scale = max(min_dim / image.width, min_dim / image.height)
        new_size = (int(image.width * scale), int(image.height * scale))
        image = image.resize(new_size, Image.LANCZOS)

    # Convert to grayscale for more stable OCR
    gray = image.convert("L")
    text = pytesseract.image_to_string(gray)
    return clean_extracted_text(text)


def _ocr_pdf_bytes(data: bytes) -> str:
    """Render PDF pages to images and OCR each page."""
    from pdf2image import convert_from_bytes

    pages = convert_from_bytes(data)
    parts: list[str] = []
    for page in pages:
        t = _ocr_image(page)
        if t:
            parts.append(t)
    return "\n\n".join(parts).strip()


def _ocr_pptx_bytes(data: bytes) -> str:
    """
    OCR image content from a PPTX file given as bytes.

    Text-based shapes are already handled by the PPT parser; here we focus on
    picture shapes (e.g. scanned documents embedded as images). As a first
    strategy, try converting the PPTX/PPT to PDF and reuse the stable PDF OCR
    path; if that fails, fall back to per-image OCR from shapes.
    """
    import io
    import os
    import subprocess
    import tempfile

    from PIL import Image
    from pptx import Presentation  # type: ignore[import]
    from pptx.shapes.group import GroupShape  # type: ignore[import]

    # 1) Try PPTX/PPT -> PDF conversion via LibreOffice (if installed), then reuse PDF OCR
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "input.pptx")
            pdf_path = os.path.join(tmpdir, "input.pdf")
            with open(pptx_path, "wb") as f:
                f.write(data)
            # libreoffice --headless --convert-to pdf --outdir <tmpdir> input.pptx
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmpdir,
                    pptx_path,
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            # LibreOffice may name the output slightly differently; pick any .pdf in tmpdir
            pdf_files = [f for f in os.listdir(tmpdir) if f.lower().endswith(".pdf")]
            if pdf_files:
                pdf_full_path = os.path.join(tmpdir, pdf_files[0])
                with open(pdf_full_path, "rb") as pf:
                    pdf_bytes = pf.read()
                pdf_text = _ocr_pdf_bytes(pdf_bytes)
                if pdf_text.strip():
                    return pdf_text
    except Exception:
        # If conversion fails for any reason, fall back to per-image OCR below
        pass

    # 2) Fallback: OCR images directly from shapes (including groups)
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    def iter_shape_images(shape) -> list:
        """Yield all picture image blobs from a shape, including within groups."""
        images = []
        image = getattr(shape, "image", None)
        if image is not None:
            images.append(image)
        # Recurse into group shapes
        if isinstance(shape, GroupShape):
            for child in shape.shapes:
                images.extend(iter_shape_images(child))
        return images

    for slide in prs.slides:
        for shape in slide.shapes:
            for image in iter_shape_images(shape):
                try:
                    with Image.open(io.BytesIO(image.blob)) as img:
                        t = _ocr_image(img)
                except Exception:
                    t = ""
                if t:
                    parts.append(t)

    return "\n\n".join(parts).strip()


def extract_text_with_ocr(
    file_content: bytes,
    file_name: str,
    content_type: Optional[str] = None,
) -> str:
    """
    Dispatch OCR based on file type.

    - PDF: render pages to images and OCR each page.
    - PPTX/PPT: OCR picture shapes.
    - Other types: currently unsupported (returns empty string).
    """
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    if (content_type and "pdf" in content_type.lower()) or ext == "pdf":
        return _ocr_pdf_bytes(file_content)
    if ext in {"pptx", "ppt"}:
        return _ocr_pptx_bytes(file_content)
    # Fallback: no OCR for this type yet
    return ""

